#!/usr/bin/env python3
"""
SPX PH Update — Template Slide Deck
Run: python3 make_template.py
Outputs: spx_ph_update_template.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Colours ──────────────────────────────────────────────────────────────────
ORANGE   = RGBColor(0xF4, 0x79, 0x20)   # SPX brand orange
ORANGE_D = RGBColor(0xC4, 0x5A, 0x08)   # darker orange for footer bars
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x1E, 0x1E, 0x1E)
LGREY    = RGBColor(0xF2, 0xF2, 0xF2)
MGREY    = RGBColor(0xB0, 0xB0, 0xB0)
GREEN    = RGBColor(0x00, 0xB0, 0x50)
RED      = RGBColor(0xC8, 0x1E, 0x1E)
LORANGE  = RGBColor(0xFF, 0xDA, 0xB4)   # light orange (alt table rows)
DARK_TXT = RGBColor(0x26, 0x26, 0x26)

W  = Inches(13.33)   # 16:9 widescreen
H  = Inches(7.5)


def prs_16x9():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


# ── Low-level helpers ─────────────────────────────────────────────────────────
def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(blank_layout)


def add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, font_size, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def add_para(tf, text, font_size, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, space_before=0):
    from pptx.oxml import parse_xml
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return p


# ── Reusable slide builders ───────────────────────────────────────────────────
HDR_H = Inches(0.72)
MARGIN = Inches(0.35)


def apply_content_chrome(slide, section_title):
    """White bg, orange header bar, section title."""
    add_rect(slide, 0, 0, W, H, LGREY)
    add_rect(slide, 0, 0, W, HDR_H, ORANGE)
    add_rect(slide, 0, HDR_H, W, Inches(0.04), ORANGE_D)
    add_textbox(slide, MARGIN, Inches(0.11), W - MARGIN*2, Inches(0.55),
                section_title, 22, bold=True, color=WHITE)


def make_divider_slide(slide, title, subtitle=""):
    """Full orange slide used for section dividers."""
    add_rect(slide, 0, 0, W, H, ORANGE)
    add_rect(slide, 0, H - Inches(0.65), W, Inches(0.65), ORANGE_D)
    add_textbox(slide, Inches(1), Inches(2.2), W - Inches(2), Inches(2.0),
                title, 44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, Inches(1), Inches(4.3), W - Inches(2), Inches(0.8),
                    subtitle, 18, color=WHITE, align=PP_ALIGN.CENTER)


def make_content_slide(slide, section_title, lines, font_size=11):
    apply_content_chrome(slide, section_title)
    add_rect(slide, MARGIN, HDR_H + Inches(0.14),
             W - MARGIN*2, H - HDR_H - Inches(0.25), WHITE)
    txb = slide.shapes.add_textbox(
        MARGIN + Inches(0.18), HDR_H + Inches(0.22),
        W - MARGIN*2 - Inches(0.36), H - HDR_H - Inches(0.45))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(font_size)
        run.font.color.rgb = DARK_TXT
        run.font.bold  = bool(line.startswith("##") or (line and not line.startswith(" ") and line.endswith(":")))


# ── OKR Table ────────────────────────────────────────────────────────────────
OKR_DATA = [
    # (pillar, metric)
    ("Maintain Cost\nAdvantage",          "Overall SPX CPO ($)"),
    ("Improve SPX\nVolume",               "% SHP ADO"),
    ("",                                  "% SPX LM Coverage / Buyer Coverage (w/ FEX)"),
    ("",                                  "% SPX LM Coverage / Buyer Coverage (w/o FEX)"),
    ("",                                  "% SPX FM Coverage / Seller Coverage"),
    ("",                                  "Cache % of platform"),
    ("Improve\nLogistics Speed",          "Average BWT SPX"),
    ("",                                  "% urban orders delivered D+1"),
    ("",                                  "% urban orders delivered D+2"),
    ("",                                  "% non-urban orders delivered D+2"),
    ("",                                  "% non-urban orders delivered D+3"),
    ("Improve SPX\nExperience / Quality", "99%ile BWT"),
    ("",                                  "% Order loss rate"),
    ("",                                  "BSC Coverage %"),
    ("Enhance Capabilities\nas 3PL",      "% Reverse Coverage"),
    ("",                                  "Non Shopee ADO"),
    ("",                                  "Non Shopee Profit per Order ($)"),
]

PILLAR_SPANS = [
    (0, 1),   # Maintain Cost Advantage — 1 row
    (1, 5),   # Improve SPX Volume — 5 rows
    (6, 5),   # Improve Logistics Speed — 5 rows
    (11, 3),  # Improve SPX Experience — 3 rows
    (14, 3),  # Enhance Capabilities — 3 rows
]


def rgb_to_hex(rgb_color):
    # RGBColor stores as hex string internally
    return str(rgb_color).upper()


def set_cell_color(cell, rgb_color):
    """Set table cell background fill via XML."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # Remove existing solidFill if any
    for existing in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(existing)
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr   = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', rgb_to_hex(rgb_color))


def set_cell_text(cell, text, font_size, bold=False, color=DARK_TXT,
                  align=PP_ALIGN.LEFT):
    cell.text = text
    tf = cell.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.size  = Pt(font_size)
            run.font.bold  = bold
            run.font.color.rgb = color


def add_okr_table(slide):
    COLS = 4
    ROWS = len(OKR_DATA) + 1   # +1 header row
    TBL_X = MARGIN
    TBL_Y = HDR_H + Inches(0.12)
    TBL_W = W - MARGIN * 2
    TBL_H = H - TBL_Y - Inches(0.38)

    col_widths = [Inches(1.75), Inches(4.85), Inches(1.85), Inches(1.85)]

    tbl = slide.shapes.add_table(ROWS, COLS, TBL_X, TBL_Y, TBL_W, TBL_H).table

    # Set column widths
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = cw

    # ── Header row ────────────────────────────────────────────────────────
    headers = ["Key Pillar", "Metrics", "[Month'YY]\nLatest / MTD", "[Month'YY]\nTarget"]
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        set_cell_color(cell, ORANGE)
        set_cell_text(cell, hdr, 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Data rows ─────────────────────────────────────────────────────────
    for ri, (pillar, metric) in enumerate(OKR_DATA):
        row_idx = ri + 1
        alt = (ri % 2 == 0)
        bg  = LGREY if alt else WHITE

        pillar_cell = tbl.cell(row_idx, 0)
        set_cell_color(pillar_cell, LORANGE if pillar else bg)
        set_cell_text(pillar_cell, pillar, 8, bold=bool(pillar), color=DARK_TXT, align=PP_ALIGN.CENTER)

        metric_cell = tbl.cell(row_idx, 1)
        set_cell_color(metric_cell, bg)
        set_cell_text(metric_cell, metric, 8.5, color=DARK_TXT)

        for ci in [2, 3]:
            c = tbl.cell(row_idx, ci)
            set_cell_color(c, bg)
            set_cell_text(c, "—", 9, color=MGREY, align=PP_ALIGN.CENTER)

    # ── Merge pillar cells ────────────────────────────────────────────────
    # python-pptx doesn't expose merge directly; use XML span attributes
    for start_row, span in PILLAR_SPANS:
        if span <= 1:
            continue
        # Set rowSpan on the anchor cell
        anchor = tbl.cell(start_row + 1, 0)
        anchor._tc.set('rowSpan', str(span))
        anchor._tc.set('vMerge', '1')
        # Mark continuation cells
        for r in range(start_row + 2, start_row + 1 + span):
            cont = tbl.cell(r, 0)
            cont._tc.set('vMerge', '0')

    return tbl


# ─────────────────────────────────────────────────────────────────────────────
# BUILD THE DECK
# ─────────────────────────────────────────────────────────────────────────────
prs = prs_16x9()

# ── Slide 1: Title ────────────────────────────────────────────────────────────
s1 = blank_slide(prs)
add_rect(s1, 0, 0, W, H, ORANGE)
add_rect(s1, 0, H - Inches(0.9), W, Inches(0.9), ORANGE_D)
add_rect(s1, 0, H - Inches(0.95), W, Inches(0.07), WHITE)
# Logo placeholder
add_rect(s1, MARGIN, Inches(0.22), Inches(1.4), Inches(0.46), WHITE)
add_textbox(s1, MARGIN + Inches(0.06), Inches(0.24), Inches(1.3), Inches(0.38),
            "SPX", 22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
# Title
add_textbox(s1, Inches(1), Inches(1.55), W - Inches(2), Inches(1.6),
            "SPX PH Update", 54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s1, Inches(1), Inches(3.2), W - Inches(2), Inches(0.75),
            "[Day Month Year]", 24, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s1, MARGIN, H - Inches(0.75), W - MARGIN*2, Inches(0.38),
            "CONFIDENTIAL", 11, color=WHITE, align=PP_ALIGN.CENTER)

# ── Slide 2: Executive Summary ────────────────────────────────────────────────
s2 = blank_slide(prs)
apply_content_chrome(s2, "Executive Summary")
add_rect(s2, MARGIN, HDR_H + Inches(0.14), W - MARGIN*2, H - HDR_H - Inches(0.25), WHITE)

txb = s2.shapes.add_textbox(
    MARGIN + Inches(0.18), HDR_H + Inches(0.2),
    W - MARGIN*2 - Inches(0.36), H - HDR_H - Inches(0.4))
tf = txb.text_frame
tf.word_wrap = True

lines_exec = [
    ("SPX remains [operationally on track / recovering from X] with [key theme 1] and [key theme 2].", 12, True),
    ("", 6, False),
    ("1.  Operations & Financial Performance (p.X–X):  [Key CPO message. Driver or implication.]", 11, False),
    ("", 4, False),
    ("2.  Capacity, Network & Expansion (p.X–X):  [Key volume/coverage message.]", 11, False),
    ("", 4, False),
    ("3.  Technology & AI Transformation (p.X–X):  [AI pilot status, accuracy %, HC impact.]", 11, False),
    ("", 4, False),
    ("4.  Loss Reduction & Automation (p.X–X):  [Loss rate trend, pilot progress.]", 11, False),
    ("", 4, False),
    ("5.  Financial Controls & Risk (p.X–X):  [COD exposure, remittance performance.]", 11, False),
]
first = True
for (text, fsize, bold) in lines_exec:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(fsize)
    run.font.bold  = bold
    run.font.color.rgb = DARK_TXT

# Status note
add_textbox(s2, MARGIN, H - Inches(0.42), W - MARGIN*2, Inches(0.32),
            "Posture framing: 'SPX remains … with …'  |  Each item must include a page ref  |  Lead with the metric, not the context",
            8, color=MGREY)

# ── Slide 3: Divider — OKR ───────────────────────────────────────────────────
s3 = blank_slide(prs)
make_divider_slide(s3, "OKR, Cost Initiatives\nand Budget Update")

# ── Slide 4: OKR Table ───────────────────────────────────────────────────────
s4 = blank_slide(prs)
apply_content_chrome(s4, "Overall OKR")
add_okr_table(s4)
# Legend
add_textbox(s4, MARGIN, H - Inches(0.32), Inches(6), Inches(0.28),
            "✅  Confident to be on track to meet target        🔴  Miss target — problem solve",
            8, color=DARK_TXT)

# ── Slide 5: P&L / CPO ───────────────────────────────────────────────────────
s5 = blank_slide(prs)
make_content_slide(s5, "P&L and CPO", [
    "*FX at [rate]",
    "",
    "[Month] P&L better by +$X.XM vs target — CPO Xc lower",
    "[Month] P&L better by +$X.XM vs target mainly from [primary driver]",
    "    +$X.Xm   [cost saving driver]",
    "    +$X.Xm   [volume uplift]",
    "    -$X.Xm   [cost increase item]",
    "",
    "CPO better by Xc with [main driver]; [secondary driver]",
    "    FM     (-0.Xc)  [reason]",
    "    SOC    (-0.Xc)  0.Xc [sub-driver 1]   /   0.Xc [sub-driver 2]",
    "    LH     (+0.Xc)  [reason]",
    "    Hub    (-0.Xc)  [reason]",
    "    LM     (+0.Xc)  [reason]",
    "    Claims (+0.Xc)  loss rate X.XX% / recovery rate XX%",
    "",
    "Compared to last month, CPO is [lower/higher] by Xc with Xc due to [reason]; balance from [other]",
], font_size=10.5)

# ── Slide 6: Speed ───────────────────────────────────────────────────────────
s6 = blank_slide(prs)
make_content_slide(s6, "Speed Updates", [
    "BWT improved to X.XX (Wk of [date]).  Urban BWT X.XX — Top X across SPX Markets.",
    "J&T Latest: X.XX  |  Internal target Mar'26: X.XX  |  Jun'26: X.XX  |  Dec'26: X.XX",
    "",
    "Metric                        Wk-4    Wk-3    Wk-2    Wk-1 (Latest)   J&T Latest   Mar'26 Target",
    "Average BWT SPX               X.XX    X.XX    X.XX    X.XX            X.XX          X.XX",
    "Average CDT SPX               X.XX    X.XX    X.XX    X.XX            X.XX          X.XX",
    "Urban CDT                     X.XX    X.XX    X.XX    X.XX            X.XX          X.XX",
    "Non-urban CDT                 X.XX    X.XX    X.XX    X.XX            X.XX          X.XX",
    "% urban orders D+1            XX.X%   XX.X%   XX.X%   XX.X%          XX.X%         XX.X%",
    "% urban orders D+2            XX.X%   XX.X%   XX.X%   XX.X%          XX.X%         XX.X%",
    "% non-urban orders D+2        XX.X%   XX.X%   XX.X%   XX.X%          XX.X%         XX.X%",
    "% non-urban orders D+3        XX.X%   XX.X%   XX.X%   XX.X%          XX.X%         XX.X%",
    "",
    "Watchpoint: [region/hub] delay driven by [cause]. Resolution by [date].",
], font_size=9.5)

# ── Slide 7: Divider — Operations ────────────────────────────────────────────
s7 = blank_slide(prs)
make_divider_slide(s7, "Operations Update")

# ── Slide 8: Last Mile ───────────────────────────────────────────────────────
s8 = blank_slide(prs)
make_content_slide(s8, "Last Mile Operations", [
    "LM productivity: X.X parcels/rider/day  (WoW: +X.X / MoM: +X.X)  vs target X.X",
    "    NCR: X.X  |  LUZ: X.X  |  VIS: X.X  |  MIN: X.X",
    "",
    "Attendance rate: XX.X%  (WoW: +X.Xpp).  Active rider pool: X,XXX",
    "",
    "Highlights:",
    "    • [Positive item — route optimisation / surge handling / high-volume hub]",
    "    • [Positive item]",
    "",
    "Watchpoints:",
    "    • [Issue + root cause + ETA to resolve]",
    "    • [Issue + region + solve in progress]",
    "",
    "Action items: [Team] to [action] by [date]",
])

# ── Slide 9: Sort Center / Middle Mile ───────────────────────────────────────
s9 = blank_slide(prs)
make_content_slide(s9, "Sort Center & Middle Mile", [
    "SOC throughput: X,XXX parcels/hr  (WoW: +XX / vs capacity: XX%)",
    "MM on-time departure: XX.X%  |  On-time arrival: XX.X%",
    "",
    "Sort accuracy: XX.XXX%  (target: XX.XXX%)",
    "",
    "Highlights:",
    "    • [Throughput gain from process change or equipment]",
    "    • [Line-haul efficiency: new route, consolidated run]",
    "",
    "Watchpoints:",
    "    • [Hub congestion / delayed departure + cause]",
    "    • [Staffing gap at [hub] — mitigation plan]",
    "",
    "Capacity utilisation: [Hub A] XX%  |  [Hub B] XX%  |  [Hub C] XX%",
])

# ── Slide 10: First Mile ─────────────────────────────────────────────────────
s10 = blank_slide(prs)
make_content_slide(s10, "First Mile", [
    "FM SLA compliance: XX.X%  (WoW: +X.Xpp / target: XX.X%)",
    "Pickup success rate: XX.X%  |  Avg pickup attempts: X.X",
    "",
    "Seller coverage (SPX FM): XX.X%  (w/ FEX: XX.X%)",
    "",
    "Highlights:",
    "    • [New FM hub or expanded coverage area]",
    "    • [Seller onboarding milestone: X new sellers activated]",
    "",
    "Watchpoints:",
    "    • [Region with low SLA compliance — root cause]",
    "    • [High no-show rate at [seller segment] — action underway]",
])

# ── Slide 11: Coverage & Expansion ───────────────────────────────────────────
s11 = blank_slide(prs)
make_content_slide(s11, "Coverage & Expansion", [
    "LM Buyer Coverage (w/ FEX): XX.X%  |  (w/o FEX): XX.X%  |  Target: XX.X%",
    "Reverse Coverage: XX.X%  |  BSC Coverage: XX.X%",
    "",
    "New branches opened MTD: X  (cumulative: XX vs full-year target: XX)",
    "Provinces with full SPX coverage: XX / XX",
    "",
    "Expansion Highlights:",
    "    • [Region / province newly activated]",
    "    • [Branch opening: location, capacity, go-live date]",
    "",
    "Watchpoints / Upcoming:",
    "    • [Delayed launch — reason + revised date]",
    "    • [Preparation for 100% SPX share ramp — decision point: Month]",
    "",
    "NSS Coverage: XX.X% of registered 3PL sellers served",
])

# ── Slide 12: Loss & Claims ───────────────────────────────────────────────────
s12 = blank_slide(prs)
make_content_slide(s12, "Loss & Claims", [
    "Order loss rate: X.XX%  (WoW: ±X.XXpp / MoM: ±X.XXpp)  vs target X.XX%",
    "Claims filed MTD: X,XXX  |  Approved: XX.X%  |  Recovery rate: XX.X%",
    "",
    "Loss breakdown by type:",
    "    Missing: XX.X%  |  Damaged: XX.X%  |  Wrong item: XX.X%  |  [Other]: XX.X%",
    "",
    "High-risk segments: [seller type / product category] — [rate] vs avg [rate]",
    "",
    "AI Loss Pilot status:",
    "    • Detection accuracy: ~XX%  |  Scope: [hubs / routes covered]",
    "    • Investigation automation: [X]% of cases auto-resolved",
    "    • Next milestone: expand to X hubs by [date]",
    "",
    "Actions: [Team] — [preventive measure] targeting [X]% reduction by [date]",
])

# ── Slide 13: NSS ─────────────────────────────────────────────────────────────
s13 = blank_slide(prs)
make_content_slide(s13, "NSS — Non-Shopee Sales", [
    "NSS ADO: X,XXX  (WoW: +XX / MoM: +XXX)  vs target X,XXX",
    "NSS Profit per Order: $0.XXX  (WoW: ±$0.XXX)",
    "",
    "Active 3PL clients: XX  |  New clients onboarded MTD: X",
    "Top verticals: [vertical 1] XX%  |  [vertical 2] XX%  |  [vertical 3] XX%",
    "",
    "Highlights:",
    "    • [New enterprise client won / RFQ closed]",
    "    • [Volume milestone or product mix improvement]",
    "",
    "Watchpoints:",
    "    • [Churn risk or SLA miss with 3PL client — action]",
    "    • [Pricing pressure in [segment]]",
    "",
    "Coverage: XX.X% of registered 3PL seller base reachable",
])

# ── Slide 14: AI & Technology ────────────────────────────────────────────────
s14 = blank_slide(prs)
make_content_slide(s14, "Technology & AI", [
    "AI Control Tower",
    "    • Delay detection accuracy: ~XX%  |  Monitoring HC reduction potential: ~XX%",
    "    • Live on: [hubs/routes]  |  Expanding to: [next scope] by [date]",
    "",
    "Hub Copilot / Route Optimisation",
    "    • Productivity uplift: +X.X parcels/rider  |  Adoption: XX% of riders",
    "    • [Feature] rolled out to X hubs this week",
    "",
    "Automation Pilots",
    "    • [Process]: [X]% tasks automated  |  Saves ~X FTE equivalent",
    "    • [Process]: In pilot at [hub] — results by [date]",
    "",
    "FinOps / COD",
    "    • COD exposure: ₱XX.XM  (MoM: ±₱X.XM)  |  Remittance D+X rate: XX.X%",
    "    • Buyer perception score: X.X / 5  (WoW: ±X.X)",
])

# ── Slide 15: Appendix divider ────────────────────────────────────────────────
s15 = blank_slide(prs)
make_divider_slide(s15, "Appendix", "Supporting data & backup slides")

# ─────────────────────────────────────────────────────────────────────────────
OUT = "/home/user/gdrive-connector/spx_ph_update_template.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
